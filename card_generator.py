import cv2 as cv
import os
from pptx import Presentation
from pptx.util import Cm
import shutil
import datetime


# 畫上圖片
def draw_on(img, img2, x, y):
    img[y:y + img2.shape[0], x: x + img2.shape[1]] = img2


# scale to fit
def scale_to_fit(img, w, h):
    img_w = img.shape[1]
    img_h = img.shape[0]

    ratio_w = w / img_w

    result_h = img_h * ratio_w
    if result_h > h:
        ratio_h = h / img_h
        result_w = img_w * ratio_h
        return cv.resize(img, (int(result_w), int(h)))
    else:
        return cv.resize(img, (int(w), int(result_h)))


# 縮放成指定面積
def scale_to_area(img, target_area):
    img_w = img.shape[1]
    img_h = img.shape[0]

    current_area = img_w*img_h
    ratio = (target_area / current_area) ** 0.5

    result_w = img_w * ratio
    result_h = img_h * ratio
    return cv.resize(img, (int(result_w), int(result_h)))


if __name__ == '__main__':
    # 定義項目
    # actions = ['睡覺', '吃飯', '玩', '跳', '走路', '投擲', '寫字', '讀',
    #            '畫', '看', '聽', '打滾', '洗澡', '剪', '貼', '唱歌', '坐',
    #            '站', '蹲', '跑']
    # actions = ['蠟筆', '印章', '貼紙', '書', '鼓', '拼圖', '鏡子', '黏土', '套圈圈', '保齡球瓶', '砧板', '刀子', '叉子', '湯匙', '鏟子', '鍋子', '盤子', '水龍頭', '瓦斯爐', '冰箱', '溜滑梯', '盪鞦韆', '積木']
    #
    # 下載項目
    # main_keywords = actions
    # supplemented_keywords = ['']

    # 下載圖片
    source = '/Users/mingtayang/Desktop/source'
    done = '/Users/mingtayang/Desktop/done'
    output_folder = '/Users/mingtayang/Desktop'

    # download_with_urllib.download(path, main_keywords, supplemented_keywords)
    # p = Pool()  # number of process is the number of cores of your CPU
    # for i in range(len(main_keywords)):
    #     p.apply_async(download_with_urllib.download_images, args=(main_keywords[i], supplemented_keywords, path))
    # p.close()
    # p.join()
    #
    # for main_keyword in main_keywords:
    #     del_unwanted.remove_corrupted_image(path+main_keyword)
    #
    # print('Finished Downloading!')
    # exit()

    # # 讀取圖片
    # # 卡通
    # # 1~17: 1 張狗, 2 張其他
    # # 18~34: 1 張狗, 1 張動物, 1 張其他
    # # 34~50: 1 張狗, 2 張動物
    #
    # # 分三區：卡通狗、非狗卡通動物、其他
    img_target = []
    # img_similars = []
    # img_unsimilars = []
    #

    extensions = ['jpg', 'png', 'jpeg', 'gif']
    unsupported_extensions = ['webp']

    for dirname, subdirnames, filenames in os.walk(source):
        for filename in filenames:
            for extension in extensions:
                if filename.endswith(extension.lower()) or filename.endswith(extension.upper()):
                    fn = os.path.join(dirname, filename)
                    img_target.append(fn)
                    continue
            for extension in unsupported_extensions:
                if filename.endswith(extension.lower()) or filename.endswith(extension.upper()):
                     raise Exception(f"Found file with unsupported extension: {extension}")

    print(img_target)

    if len(img_target) == 0:
        raise Exception("沒有圖片！！")

    source_images = img_target.copy()

    rows = []
    row = []
    for _ in range(len(img_target)):
        row.append(img_target.pop())
        if len(row) == 3:
            rows.append(row)
            row = []

    if len(row):
        rows.append(row)

    # 縮放圖片
    area = Cm(6).pt * Cm(6).pt

    # 畫上圖片
    prs = Presentation()
    prs.slide_width = Cm(29.7)
    prs.slide_height = Cm(21)
    blank_slide_layout = prs.slide_layouts[6]

    upper_padding = Cm(0)
    lower_padding = Cm(0)
    col_spacing = Cm(9.9)
    left_x = Cm(0)

    imagelist = []
    row_count = len(rows)
    r = row_count//2 if row_count%2 == 0 else row_count//2+1
    for i in range(r): # 0, 1, ... , 24
        slide = prs.slides.add_slide(blank_slide_layout)

        upper_images = rows[i*2]
        lower_images = rows[i*2+1] if i*2+1 < len(rows) else None

        upper_row_padding = 40
        lower_row_padding = 40

        # 評分 padding
        # x_upper = upper_row_padding
        x = left_x
        for index, image in enumerate(upper_images):
            im = cv.imread(image, cv.IMREAD_IGNORE_ORIENTATION)

            resized = scale_to_area(im, area)
            width = Cm(resized.shape[1] / Cm(1).pt)
            height = Cm(resized.shape[0] / Cm(1).pt)

            # x = Cm(x_upper / Cm(1).pt)
            y = upper_padding

            # 最後一張靠右
            if index == len(upper_images) - 1:
                x = prs.slide_width - width

            img = slide.shapes.add_picture(image, x, y, width=width, height=height)
            # x_upper += resized.shape[1] + upper_row_padding
            x += col_spacing

        if lower_images:
            x = left_x
            # x_lower = lower_row_padding
            for index, image in enumerate(lower_images):
                im = cv.imread(image, cv.IMREAD_IGNORE_ORIENTATION)

                resized = scale_to_area(im, area)
                width = Cm(resized.shape[1] / Cm(1).pt)
                height = Cm(resized.shape[0] / Cm(1).pt)

                # x = Cm(x_lower / Cm(1).pt)
                y = Cm(prs.slide_height.cm - height.cm - lower_padding.cm)

                # 最後一張靠右
                if index == len(upper_images) - 1:
                    x = prs.slide_width - width

                img = slide.shapes.add_picture(image, x, y, width=width, height=height)
                # x_lower += resized.shape[1] + lower_row_padding
                x += col_spacing

    now = datetime.datetime.now().strftime("%m%d-%H%M")

    prs.save(f"{output_folder}/{now}.pptx")  # saving file

    # 移動完成 source 到 done 資料夾
    # for image in source_images:
    #     target = f"{done}/{os.path.basename(image)}"
    #     shutil.move(image, target)
