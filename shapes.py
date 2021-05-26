# import cv2 as cv
# import numpy as np
#
# # 建立 A4 圖片
# # A4 2479×3508
# total_width = 3508
# total_height = 2479
# # 全黑(數值為 0)
# canvas = np.zeros((total_height, total_width, 3), np.uint8)
#
# # 統一改顏色
# canvas[:] = (255, 255, 255)
#
# # cv.circle(canvas,(100, 0), 25, (0,255,0))
# #
# # cv.circle(canvas,(0, 100), 25, (0,0,255))
# #
# # cv.circle(canvas,(100, 100), 50, (255,0,0), 3)
#
# color = (255,0,0)
# a = int(96/2.54*6)
# padding = a + 250
# cv.circle(canvas, (250, 250), a, color, -1)
# cv.circle(canvas, (250+padding, 250), a, color, -1)
# cv.circle(canvas, (250+padding+padding, 250), a, color, -1)
# cv.circle(canvas, (250, 250+padding), a, color, -1)
# cv.circle(canvas, (250+padding, 250+padding), a, color, -1)
# cv.circle(canvas, (250+padding+padding, 250+padding), a, color, -1)
#
# cv.imwrite('test.jpg', canvas)

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import datetime


def make_gray_dash(shape, remove_fill=False):
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.color.brightness = 0.3

    if remove_fill:
        shape.fill.background()


output_folder = '/Users/mingtayang/Desktop'

prs = Presentation()
width = 29.7
height = 21
prs.slide_width = Cm(width)
prs.slide_height = Cm(height)
blank_slide_layout = prs.slide_layouts[6]

left_padding = 2
right_padding = 2
upper_padding = 2
lower_padding = 2
object_padding = 4
# col_spacing = Cm(9.9)
# left_x = Cm(0)

# imagelist = []
# row_count = len(rows)
# r = row_count // 2 if row_count % 2 == 0 else row_count // 2 + 1
# for i in range(r):  # 0, 1, ... , 24

slide = prs.slides.add_slide(blank_slide_layout)

# shapes = slide.shapes
# left = top = width = height = Cm(1.0)
# shape = shapes.add_shape(
#     MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
# )
#
# line = shape.line
# line.color.rgb = RGBColor(0, 0, 0)
# line.color.brightness = 0.5  # 50% lighter
# line.width = Cm(2.5)

line_length = 21

# straight
# for i in range(int(width//object_padding)):
#     line1=slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Cm(left_padding+object_padding*i), Cm(upper_padding), Cm(left_padding+object_padding*i), Cm(height-lower_padding))
#     make_gray_dash(line1)

# cross
# 上
# cross_diff = 5
# begin_x = left_padding
# end_y = None
# while True:
#     end_x = begin_x + cross_diff
#     if end_x > width - left_padding:
#         break
#
#     begin_y = upper_padding
#     end_y = begin_y + cross_diff
#
#     line1=slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Cm(begin_x), Cm(begin_y), Cm(end_x), Cm(end_y))
#     make_gray_dash(line1)
#
#     line2 = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Cm(begin_x), Cm(end_y), Cm(end_x), Cm(begin_y))
#     make_gray_dash(line2)
#
#     begin_x = end_x + object_padding
#
# # 下
# begin_x = left_padding
# begin_y = end_y + object_padding
# print(end_y)
# print(begin_y)
# while True:
#     end_x = begin_x + cross_diff
#     if end_x > width - left_padding:
#         break
#
#     end_y = begin_y + cross_diff
#
#     line1=slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Cm(begin_x), Cm(begin_y), Cm(end_x), Cm(end_y))
#     make_gray_dash(line1)
#
#     line2 = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Cm(begin_x), Cm(end_y), Cm(end_x), Cm(begin_y))
#     make_gray_dash(line2)
#
#     begin_x = end_x + object_padding


# diagonal
# 上
# cross_diff = 5
# begin_x = left_padding
# end_y = None
# while True:
#     end_x = begin_x + cross_diff
#     if end_x > width - left_padding:
#         break
#
#     begin_y = upper_padding
#     end_y = begin_y + cross_diff
#
#     line1=slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Cm(begin_x), Cm(begin_y), Cm(end_x), Cm(end_y))
#     make_gray_dash(line1)
#
#     begin_x = end_x + object_padding/3
#
# # 下
# begin_x = left_padding
# begin_y = end_y + object_padding
# print(end_y)
# print(begin_y)
# while True:
#     end_x = begin_x + cross_diff
#     if end_x > width - left_padding:
#         break
#
#     end_y = begin_y + cross_diff
#
#     line2 = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Cm(begin_x), Cm(end_y), Cm(end_x), Cm(begin_y))
#     make_gray_dash(line2)
#
#     begin_x = end_x + object_padding/3

# circle, rectangle, triangle
# my_shape = MSO_AUTO_SHAPE_TYPE.OVAL
# my_shape = MSO_AUTO_SHAPE_TYPE.RECTANGLE
my_shape = MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE

# 上
diameter = 5
begin_x = left_padding
end_y = None
while True:
    end_x = begin_x + diameter
    if end_x > width - left_padding:
        break

    begin_y = upper_padding
    end_y = begin_y + diameter

    added_shape = slide.shapes.add_shape(my_shape, Cm(begin_x), Cm(begin_y), Cm(end_x-begin_x), Cm(end_y-begin_y))
    make_gray_dash(added_shape, remove_fill=True)

    begin_x = end_x + object_padding

# 下
begin_x = left_padding
begin_y = end_y + object_padding
while True:
    end_x = begin_x + diameter
    if end_x > width - left_padding:
        break

    end_y = begin_y + diameter

    added_shape = slide.shapes.add_shape(my_shape, Cm(begin_x), Cm(begin_y), Cm(end_x-begin_x), Cm(end_y-begin_y))
    make_gray_dash(added_shape, remove_fill=True)

    begin_x = end_x + object_padding

# 產出
now = datetime.datetime.now().strftime("%m%d-%H%M")
prs.save(f"{output_folder}/{now}.pptx")  # saving file
