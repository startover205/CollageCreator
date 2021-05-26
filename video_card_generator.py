import cv2 as cv
import numpy as np
import os
import random
import copy
from fpdf import FPDF
import download_with_urllib
from multiprocessing import Pool
import del_unwanted
from pptx import Presentation
from pptx.util import Cm


# 畫上圖片
def draw_on(img, img2, x, y):
    img[y:y + img2.shape[0], x: x + img2.shape[1]] = img2


# scale to fit
def scale_to_fit(img, w, h):
    img_w = img.shape[1]
    img_h = img.shape[0]

    ratio_w = w / img_w

    result_h = img_h * ratio_w
    if result_h > h:
        ratio_h = h / img_h
        result_w = img_w * ratio_h
        return cv.resize(img, (int(result_w), int(h)))
    else:
        return cv.resize(img, (int(w), int(result_h)))


# 縮放成指定面積
def scale_to_area(img, target_area):
    img_w = img.shape[1]
    img_h = img.shape[0]

    current_area = img_w*img_h
    ratio = (target_area / current_area) ** 0.5

    result_w = img_w * ratio
    result_h = img_h * ratio
    return cv.resize(img, (int(result_w), int(result_h)))


if __name__ == '__main__':
    # 定義項目
    # actions = ['睡覺', '吃飯', '玩', '跳', '走路', '投擲', '寫字', '讀',
    #            '畫', '看', '聽', '打滾', '洗澡', '剪', '貼', '唱歌', '坐',
    #            '站', '蹲', '跑']
    actions = ['蠟筆', '印章', '貼紙', '書', '鼓', '拼圖', '鏡子', '黏土', '套圈圈', '保齡球瓶', '砧板', '刀子', '叉子', '湯匙', '鏟子', '鍋子', '盤子', '水龍頭', '瓦斯爐', '冰箱', '溜滑梯', '盪鞦韆', '積木']
    #
    # 下載項目
    main_keywords = actions
    supplemented_keywords = ['']

    # 下載圖片
    path = './data/actions2/'
    # download_with_urllib.download(path, main_keywords, supplemented_keywords)
    # p = Pool()  # number of process is the number of cores of your CPU
    # for i in range(len(main_keywords)):
    #     p.apply_async(download_with_urllib.download_images, args=(main_keywords[i], supplemented_keywords, path))
    # p.close()
    # p.join()
    #
    # for main_keyword in main_keywords:
    #     del_unwanted.remove_corrupted_image(path+main_keyword)
    #
    # print('Finished Downloading!')
    # exit()

    # # 讀取圖片
    # # 卡通
    # # 1~17: 1 張狗, 2 張其他
    # # 18~34: 1 張狗, 1 張動物, 1 張其他
    # # 34~50: 1 張狗, 2 張動物
    #
    # # 分三區：卡通狗、非狗卡通動物、其他
    img_target = []
    # img_similars = []
    # img_unsimilars = []
    #
    for dirname, subdirnames, filenames in os.walk(path):
        for filename in filenames:
            if filename.endswith('mp4'):
                fn = os.path.join(dirname, filename)
                # if os.path.basename(dirname) in target:
                img_target.append(fn)
                # elif os.path.basename(dirname) in similars:
                #     img_similars.append(fn)
                # else:
                #     img_unsimilars.append(fn)
    #
    # print(similars)
    #
    # random.shuffle(img_target)
    # random.shuffle(img_similars)
    # random.shuffle(img_unsimilars)
    # print(len(img_target), len(img_similars), len(img_unsimilars))
    #
    # # 分類成列
    # # 1~17: 1 張狗, 2 張其他
    # # 18~34: 1 張狗, 1 張動物, 1 張其他
    # # 34~50: 1 張狗, 2 張動物
    #
    rows = []
    row = []
    print(len(img_target))
    for _ in range(len(img_target)):
        row.append(img_target.pop())
        if len(row) == 2:
            rows.append(row)
            row = []

    if len(row):
        rows.append(row)

    print(rows)
    # 排列圖列拼貼

    # 畫上圖片
    prs = Presentation()
    prs.slide_width = Cm(29.7)
    prs.slide_height = Cm(21)
    blank_slide_layout = prs.slide_layouts[6]

    upper_padding = Cm(1)
    lower_padding = Cm(1)
    col_spacing = Cm(1)

    imagelist = []
    row_count = len(rows)
    r = row_count//2 if row_count%2 == 0 else row_count//2+1
    for i in range(r): # 0, 1, ... , 24
        slide = prs.slides.add_slide(blank_slide_layout)

        upper_images = rows[i*2]
        lower_images = rows[i*2+1] if i*2+1 < len(rows) else None

        print(upper_images)
        print(lower_images)

        # 評分 padding
        # x_upper = upper_row_padding
        x = col_spacing
        for index, image in enumerate(upper_images):
            width = Cm(13.35)
            height = Cm(7.5)

            # x = Cm(x_upper / Cm(1).pt)
            y = upper_padding

            # img = slide.shapes.add_picture(image, x, y, width=width, height=height)
            img = slide.shapes.add_movie(image, x, y, width=width, height=height)
            x += width + col_spacing

        if lower_images:
            x = col_spacing
            # x_lower = lower_row_padding
            for index, image in enumerate(lower_images):
                width = Cm(13.35)
                height = Cm(7.5)

                # x = Cm(x_lower / Cm(1).pt)
                y = Cm(prs.slide_height.cm - height.cm - lower_padding.cm)

                img = slide.shapes.add_movie(image, x, y, width=width, height=height)
                x += width + col_spacing

    prs.save("slide3.pptx")  # saving file
