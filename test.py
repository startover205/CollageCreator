import cv2 as cv
import os
from pptx import Presentation
from pptx.util import Cm
import datetime


# 畫上圖片
def draw_on(img, img2, x, y):
    img[y:y + img2.shape[0], x: x + img2.shape[1]] = img2


# scale to fit
def scale_to_fit(img, w, h):
    img_w = img.shape[1]
    img_h = img.shape[0]

    ratio_w = w / img_w

    result_h = img_h * ratio_w
    if result_h > h:
        ratio_h = h / img_h
        result_w = img_w * ratio_h
        return cv.resize(img, (int(result_w), int(h)))
    else:
        return cv.resize(img, (int(w), int(result_h)))


# 縮放成指定面積
def scale_to_area(img, target_area):
    img_w = img.shape[1]
    img_h = img.shape[0]

    current_area = img_w*img_h
    ratio = (target_area / current_area) ** 0.5

    result_w = img_w * ratio
    result_h = img_h * ratio
    return cv.resize(img, (int(result_w), int(result_h)))


def cm_width(img):
    return Cm(img.shape[1] / Cm(1).pt)


def cm_height(img):
    return Cm(img.shape[0] / Cm(1).pt)


def cm(pt):
    return  Cm(pt/ Cm(1).pt)


if __name__ == '__main__':
    # 定義項目
    # actions = ['睡覺', '吃飯', '玩', '跳', '走路', '投擲', '寫字', '讀',
    #            '畫', '看', '聽', '打滾', '洗澡', '剪', '貼', '唱歌', '坐',
    #            '站', '蹲', '跑']
    # actions = ['蠟筆', '印章', '貼紙', '書', '鼓', '拼圖', '鏡子', '黏土', '套圈圈', '保齡球瓶', '砧板', '刀子', '叉子', '湯匙', '鏟子', '鍋子', '盤子', '水龍頭', '瓦斯爐', '冰箱', '溜滑梯', '盪鞦韆', '積木']
    #
    # 下載項目
    # main_keywords = actions
    # supplemented_keywords = ['']

    # 下載圖片
    source = '/Users/mingtayang/Desktop/source'
    done = '/Users/mingtayang/Desktop/done'
    output_folder = '/Users/mingtayang/Desktop'

    # download_with_urllib.download(path, main_keywords, supplemented_keywords)
    # p = Pool()  # number of process is the number of cores of your CPU
    # for i in range(len(main_keywords)):
    #     p.apply_async(download_with_urllib.download_images, args=(main_keywords[i], supplemented_keywords, path))
    # p.close()
    # p.join()
    #
    # for main_keyword in main_keywords:
    #     del_unwanted.remove_corrupted_image(path+main_keyword)
    #
    # print('Finished Downloading!')
    # exit()

    # # 讀取圖片
    # # 卡通
    # # 1~17: 1 張狗, 2 張其他
    # # 18~34: 1 張狗, 1 張動物, 1 張其他
    # # 34~50: 1 張狗, 2 張動物
    #
    # # 分三區：卡通狗、非狗卡通動物、其他
    img_target = []
    # img_similars = []
    # img_unsimilars = []
    #

    extensions = ['jpg', 'png', 'jpeg', 'gif']
    unsupported_extensions = ['webp']

    for dirname, subdirnames, filenames in os.walk(source):
        for filename in filenames:
            for extension in extensions:
                if filename.endswith(extension.lower()) or filename.endswith(extension.upper()):
                    fn = os.path.join(dirname, filename)
                    img_target.append(fn)
                    continue
            for extension in unsupported_extensions:
                if filename.endswith(extension.lower()) or filename.endswith(extension.upper()):
                    raise Exception(f"Found file with unsupported extension: {extension}")

    print(img_target)

    if len(img_target) == 0:
        raise Exception("沒有圖片！！")

    source_images = img_target.copy()

    # rows = []
    # row = []
    # for _ in range(len(img_target)):
    #     row.append(img_target.pop())
    #     if len(row) == 3:
    #         rows.append(row)
    #         row = []

    # if len(row):
    #     rows.append(row)

    # 縮放圖片
    area = Cm(6).pt * Cm(6).pt

    # 畫上圖片
    prs = Presentation()
    prs.slide_width = Cm(29.7)
    prs.slide_height = Cm(21)
    blank_slide_layout = prs.slide_layouts[6]

    # 算出所有圖片的寬高比、按照大小排列、統一轉成垂直圖片

    # 將紙張橫向放置，讓圖卡由寬高比低開始排列，有上而下，超過高度邊界則換行，再繼續由下而上，反覆來回，直到超過寬度邊界再換紙
    temp_fn = 'temp.jpg'

    items = []
    for image in img_target:
        im = cv.imread(image, cv.IMREAD_IGNORE_ORIENTATION+cv.IMREAD_COLOR)

        resized = scale_to_area(im, area)
        width = resized.shape[1]
        height = resized.shape[0]
        if width > height:
            resized = cv.rotate(resized, cv.ROTATE_90_CLOCKWISE)
            im = cv.rotate(im, cv.ROTATE_90_CLOCKWISE)
            items.append((im, height, width))
        else :
            items.append((im, width, height))


            # resized = rotate(resized, 90)

        # images.append(im, width, height)

    # sort by width
    items = sorted(items, key=lambda item: item[1])

    # for image in imagelist:
    #     print(image.shape[1])

    last_width = Cm(0)
    max_height = None
    top_margin = left_margin = bottom_margin = right_margin = Cm(1)
    horizontal_padding = vertical_padding = Cm(1)
    # col_spacing = Cm(9.9)
    # left_x = Cm(0)

    # last_height = Cm(0)
    x = left_margin
    y = top_margin
    slide = prs.slides.add_slide(blank_slide_layout)
    for item in items:
        width = cm(item[1])
        height = cm(item[2])

        # 確認不會超出紙張高度，超出就換列
        if y + height + bottom_margin > prs.slide_height:
            y = top_margin
            x += horizontal_padding
            x += last_width

            # 如果換行後，超出紙張寬度，則建立新的紙張
            if x + width + right_margin > prs.slide_width:
                slide = prs.slides.add_slide(blank_slide_layout)
                x = left_margin
                max_height = None

        # 紀錄此張投影片最大高度圖片
        if max_height is None:
            max_height = height

        cv.imwrite(temp_fn, item[0])
        a = slide.shapes.add_picture(temp_fn, x, y, width=width, height=height)
        last_width = width
        # last_height = height

        y += max_height
        y += vertical_padding

    # 刪除暫存
    os.remove("temp.jpg")

    # 輸出 ppt
    now = datetime.datetime.now().strftime("%m%d-%H%M")

    prs.save(f"{output_folder}/{now}.pptx")  # saving file

    # 移動完成 source 到 done 資料夾
    # for image in source_images:
    #     target = f"{done}/{os.path.rbasename(image)}"
    #     shutil.move(image, target)
